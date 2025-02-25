from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from dotenv import load_dotenv
from langchain_chroma import Chroma
from langchain_upstage import UpstageEmbeddings
from langchain_core.documents import Document
import re, os, json

load_dotenv()

def extract_medical_codes_from_pptx(pptx_path, pediatric_start_page=None):
      # 파일 존재 확인
    if not os.path.exists(pptx_path):
        print(f"파일을 찾을 수 없습니다: {pptx_path}")
        print("현재 디렉토리:", os.getcwd())
        print("디렉토리 내 파일:", os.listdir())
        raise FileNotFoundError(f"파일을 찾을 수 없습니다: {pptx_path}")

    prs = Presentation(pptx_path)

    if pediatric_start_page is not None:
        print(f"소아 슬라이드 {pediatric_start_page}로 지정됨.")

    # 구조화된 데이터 형식
    organized_data = {}

    # 중복 방지
    processed_items = set()

    for idx, slides in enumerate(prs.slides, 1):
        # if idx < 191 or idx > 195:
        #     continue

        print(f"\n=== 슬라이드 {idx} ===")
        print(f"도형 수: {len(slides.shapes)}")   # 각 슬라이드의 도형수
        nacrs_code = None
        title = None
        current_category = None

        # 소아 시작 페이지가 있고, 그 페이지 이상 페이지라면 소아 페이지 지정
        is_pediatric = (pediatric_start_page is not None and idx >= pediatric_start_page)
        
        # 슬라이드의 모든 도형에서 텍스트 추출
        for shape in slides.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                print("표 발견")
                table = shape.table

                # 표 처리 
                for row in table.rows:
                    row_text = " ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    print(f"처리할 행: {row_text}")  # 디버깅용

                    # 제목 찾기(대분류)
                    if 'Coding System' in row_text:
                        parts = row_text.split('Codes')
                        if len(parts) > 1:
                            title = parts[1].strip()
                            print(f"제목: {title}")  # 디버깅용

                    # NACRS 코드 먼저 찾기
                    elif 'NACRS' in row_text:
                        code_match = re.search(r'NACRS\s+(\d+)', row_text)
                        if code_match:
                            nacrs_code = code_match.group(1)
                            print(f"NACRS 코드: {nacrs_code}")  # 디버깅용

                        # 아직 나오지 않은 코드라면
                        if nacrs_code and nacrs_code not in organized_data:
                            organized_data[nacrs_code] = {
                                'title' : title,
                                'adult':{
                                    'vital_signs_primary':{},
                                    'other_primary':{},
                                    'symptom_secondary':{}
                                },
                                'pediatric' : {
                                    'vital_signs_primary':{},
                                    'other_primary':{},
                                    'symptom_secondary':{}
                                }
                            }
                    
                    # 카테고리 확인
                    if '활력징후 1차 고려사항' in row_text:
                        current_category = 'vital_signs_primary'
                    elif '그 밖의 1차 고려사항' in row_text:
                        current_category = 'other_primary'
                    elif '증상별 2차 고려사항' in row_text:
                        current_category = 'symptom_secondary'
                    else:
                        # 레벨과 설명 매칭
                        level_matches = re.finditer(r'(\d+)\s+(.*?)(?=\s+\d+\s+|$)', row_text)
                        for match in level_matches:
                            # 현재 코드가 있고 카테고리가 있다면
                            if nacrs_code and current_category:
                                level = match.group(1)
                                desc = match.group(2).strip()
                                if desc:
                                    # 성인/소아 구분하여 저장
                                    patient_type = 'pediatric' if is_pediatric else 'adult'

                                    # 중복 방지 위해 고유 키 생성
                                    item_key = f"{nacrs_code}_{patient_type}_{current_category}_{level}_{desc}"
                                    if item_key not in processed_items:
                                        # 같은 레벨끼리 묶기 위해 고유 키 사용
                                        if level not in organized_data[nacrs_code][patient_type][current_category]:
                                            organized_data[nacrs_code][patient_type][current_category][level] = []

                                        # 해당 레에 항목 추가
                                        organized_data[nacrs_code][patient_type][current_category][level].append(desc)
                                        processed_items.add(item_key)
                                        print(f"항목 추가: {current_category} - 레벨 {level}, 설명 {desc}")
                                    else:
                                        print(f"중복 항목 건너뜀: {patient_type} - {current_category} - 레벨 {level}, 설명 {desc}")

    return organized_data


# 문서 변환
def convert_to_documents(data):
    """
    추출된 의학 코드 데이터를 LangChain Document 객체 리스트로 변환
    """
    documents = []
    for code, code_data in data.items():
        title = code_data.get('title')
        
        # 성인 데이터
        for category_name, category_data in code_data['adult'].items():
            for level, description in category_data.items():
                for desc in description:
                    # 각 항목을 문서로 변환
                    content = f"NACRS 코드: {code}\n"
                    content = f"제목: {title}\n"
                    content = f"환자 유형: 성인\n"
                    content = f"카테고리: {category_name}\n"
                    content = f"레벨: {level}\n"
                    content = f"설명: {desc}"

                    # 메타데이터
                    metadata = {
                        "code": code,
                        "title": title,
                        "patient_type": "adult",
                        "category": category_name,
                        "level": level,
                    }
                    
                    documents.append(Document(page_content=content, metadata=metadata))
                    
        # 소아 데이터 처리
        for categroy_name, category_data in code_data['pediatric'].items():
            for level, description in category_data.items():
                for desc in description:
                    # 각 항목을 문서로 변환
                    # 각 항목을 문서로 변환
                    content = f"NACRS 코드: {code}\n"
                    content = f"제목: {title}\n"
                    content = f"환자 유형: 소아\n"
                    content = f"카테고리: {category_name}\n"
                    content = f"레벨: {level}\n"
                    content = f"설명: {desc}"
                    
                    # 메타데이터
                    metadata = {
                        "code": code,
                        "title": title,
                        "patient_type": "pediatric",
                        "category": category_name,
                        "level": level,
                    }
            
                    documents.append(Document(page_content=content, metadata=metadata))
                    
    print(f"{len(documents)}개의 문서로 변환됨.")
    return documents
    
# vectorStore
def get_vectorstore(pptx_path=None, pediatric_start_page=None):
    # chroma db 경로 설정
    persist_directory = "./chroma_db"

    # embeddings
    embeddings = UpstageEmbeddings(
        api_key=os.getenv("UPSTAGE_API_KEY"),
        model="embedding-passage"
    )
    
    # 중복 방지
    if os.path.exists(persist_directory):
        print("기존 db 사용")
        return Chroma(
            persist_directory=persist_directory,
            embedding_function=embeddings
        )
    else:
        print("새로운 db 생성")
        # 문서 준비 및 Vector Store 생성
        if pptx_path is None:
            pptx_path = "src/KTAS_guideline.pptx"
        
        # 데이터 추출
        data = extract_medical_codes_from_pptx(pptx_path, pediatric_start_page)
        
        # 백업용
        output_path = "의학코드_추출결과.json"
        save_to_json(data, output_path)
        
        # 문서 변환
        documents = convert_to_documents(data)
        
        return Chroma.from_documents(
            documents=documents,
            embedding=embeddings,
            persist_directory=persist_directory
        )



# json 변환
def save_to_json(data, output_path):
    # JSON 파일로 저장 (한글 인코딩 처리)
    with open(output_path, 'w', encoding='utf-8') as f:
        json.dump(data, f, ensure_ascii=False, indent=2)



# 사용 예시
if __name__ == "__main__":
    pptx_path = "scr/KTAS_guideline.pptx"   # PPT 파일 경로 
    output_path = "의학코드_추출결과.json"  # 결과 저장할 JSON 파일 경로
    pediatric_start_page = 192

    try:
        # 데이터 추출
        data = extract_medical_codes_from_pptx(pptx_path, pediatric_start_page)
        
        # json 저장
        save_to_json(data, output_path)
        
        # 벡터 스토어 생성
        vectorstore = get_vectorstore(pptx_path, pediatric_start_page)
        
    except Exception as e:
        print(f"오류 발생: {str(e)}")


